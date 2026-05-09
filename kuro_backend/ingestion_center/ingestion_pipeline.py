from __future__ import annotations

import json
import logging
from pathlib import Path
from typing import Any, Dict, List

from kuro_backend import observability
from kuro_backend.ingestion_center import ingestion_audit

from . import ingestion_registry, semantic_registry
from .chunking_engine import clean_text, semantic_chunk
from .embedding_manager import rebuild_vectors

logger = logging.getLogger(__name__)
logger.propagate = False


def parse_file(file_path: str, source_type: str) -> Dict[str, Any]:
    path = Path(file_path)
    ext = (source_type or path.suffix.lstrip(".")).lower()
    with observability.trace_node("ingestion_parse", {"source_type": ext, "file_path": str(path)}):
        if ext in {"txt", "md"}:
            return {"text": path.read_text(encoding="utf-8", errors="replace"), "parser_type": ext}
        if ext == "pdf":
            import pdfplumber

            parts: List[str] = []
            with pdfplumber.open(str(path)) as pdf:
                for page in pdf.pages:
                    parts.append(page.extract_text() or "")
            return {"text": "\n\n".join(parts), "parser_type": "pdfplumber"}
        if ext == "docx":
            from docx import Document

            doc = Document(str(path))
            return {"text": "\n".join(p.text for p in doc.paragraphs), "parser_type": "python-docx"}
        if ext in {"xlsx", "xls"}:
            from openpyxl import load_workbook

            book = load_workbook(filename=str(path), data_only=True)
            rows = []
            for sheet in book.worksheets:
                rows.append(f"# Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    rows.append(" | ".join("" if value is None else str(value) for value in row))
            return {"text": "\n".join(rows), "parser_type": "openpyxl"}
        if ext in {"pptx", "ppt"}:
            from pptx import Presentation

            presentation = Presentation(str(path))
            slides = []
            for idx, slide in enumerate(presentation.slides, start=1):
                slides.append(f"# Slide {idx}")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        slides.append(text)
            return {"text": "\n".join(slides), "parser_type": "python-pptx"}
        raise ValueError(f"Unsupported parser source_type={ext}")


def _build_summary(chunks: List[Dict[str, Any]]) -> str:
    previews = [chunk["preview_text"] for chunk in chunks[:3] if chunk.get("preview_text")]
    return " ".join(previews)[:600]


def run_ingestion(dataset_uuid: str, username: str, job_id: int | None = None) -> Dict[str, Any]:
    dataset = ingestion_registry.get_dataset(dataset_uuid)
    if dataset is None:
        raise ValueError(f"Unknown dataset_uuid={dataset_uuid}")
    if job_id:
        ingestion_registry.update_job(job_id, status="processing", progress_percent=5, started_at=ingestion_registry.now_iso())
    ingestion_registry.update_dataset(dataset_uuid, ingestion_status="processing", last_error=None)
    parsed = parse_file(dataset["file_path"], dataset.get("source_type") or "")
    text = clean_text(parsed["text"])
    chunks = semantic_chunk(text, dataset_uuid)
    if not chunks:
        raise ValueError("Parsed file produced no ingestible text.")
    if job_id:
        ingestion_registry.append_job_log(job_id, "Parsed file and generated chunks.")
        ingestion_registry.update_job(job_id, progress_percent=40)
    with observability.trace_node("ingestion_embed", {"dataset_uuid": dataset_uuid, "chunk_count": len(chunks)}):
        vector_result = rebuild_vectors(dataset_uuid, chunks, dataset)
    registered = semantic_registry.register_chunks(dataset_uuid, chunks, vector_result)
    entity_count = sum(len(chunk.get("entities", [])) for chunk in chunks)
    status = "completed" if vector_result.get("status") == "completed" else "partially_indexed"
    updated = ingestion_registry.update_dataset(
        dataset_uuid,
        ingestion_status=status,
        chunk_count=registered["chunk_count"],
        embedding_count=vector_result.get("embedding_count", 0),
        vector_collection=vector_result.get("collection_name", ""),
        memory_scope=vector_result.get("memory_scope", "chroma_only"),
        parser_type=parsed.get("parser_type"),
        entity_count=entity_count,
        summary_text=_build_summary(chunks),
        metadata_json=json.dumps({"text_length": len(text)}, ensure_ascii=False),
        last_error=vector_result.get("error"),
    )
    if job_id:
        ingestion_registry.update_job(
            job_id,
            status=status,
            progress_percent=100,
            completed_at=ingestion_registry.now_iso(),
            error_message=vector_result.get("error"),
        )
    ingestion_audit.log_lineage(
        dataset_uuid,
        "ingest",
        {"chunk_count": registered["chunk_count"], "embedding_count": vector_result.get("embedding_count", 0), "status": status},
    )
    return {"dataset": updated, "job": ingestion_registry.get_job(job_id) if job_id else None, "chunks": registered["chunk_count"]}


def run_reindex(dataset_uuid: str, username: str, job_id: int | None = None) -> Dict[str, Any]:
    result = run_ingestion(dataset_uuid, username, job_id=job_id)
    ingestion_audit.log_lineage(dataset_uuid, "reindex", {"status": result["dataset"]["ingestion_status"]})
    return result
