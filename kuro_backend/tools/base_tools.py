"""
Kuro AI V6.0 Sovereign - Tools [2026-04-17]
============================================================
Kuro OS Sentinel - Universal File Parser & System Analyzer
Butler + System Administrator capabilities.
Supports: Text, Code, PDF, Images (Vision), Logs, and recursive crawling.

PHASE 1 Fixes [2026-04-05]:
- Path Integrity: All file interactions use os.path.abspath(PROJECT_ROOT)
- Physical Validation: os.path.exists() checks before file operations

--- Header Doc ---
Purpose: All Gemini-callable tools (filesystem, system inspection, reminders, habits, finance, market, OpenClaw advanced execution).
Caller: core.py, langgraph_core tool_node, services/core_service (some delegated helpers), tests.
Dependencies: reminder_service, finance_db, execution.openclaw_bridge, chromadb (context lookup), PDF/DOCX parsers, requests, threading.
Main Functions: get_system_status, universal_read, smart_read, analyze_compliance, add_reminder_tool, set_monthly_budget_tool, get_ticker_price_tool, get_market_news_tool, prediction_market_scan_tool, advanced_execution_tool.
Side Effects: Filesystem reads under PROJECT_ROOT whitelist, subprocess calls, finance_db writes for price/prediction cache, OpenClaw HTTP calls via bridge.
"""
import asyncio
import base64
import json
import logging
import os
import re
import gzip
import threading
import urllib3
import requests
import psutil
import subprocess
from datetime import datetime, timedelta
from typing import Any, Dict, List, Optional, Tuple
from kuro_backend.config import settings

logger = logging.getLogger(__name__)
logger.propagate = False  # Prevent double-reporting to root logger

# Disable SSL warnings for self-signed Proxmox certificates
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# --- Configuration ---
# CRITICAL: Define PROJECT_ROOT as absolute path to avoid path confusion
PROJECT_ROOT = os.path.abspath(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
UPLOAD_DIR = os.path.abspath(os.path.join(PROJECT_ROOT, "uploaded_files"))
LOGS_DIR = os.path.abspath(os.path.join(PROJECT_ROOT, "logs"))
DB_DIR = os.path.abspath(os.path.join(PROJECT_ROOT, "db"))

# Ensure directories exist
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(LOGS_DIR, exist_ok=True)
os.makedirs(DB_DIR, exist_ok=True)

MAX_FILE_SIZE_MB = 50  # Skip files larger than 50MB
WHITELIST_PATHS = [
    "/home/kuro/projects/",
    "/var/log/",
    UPLOAD_DIR,  # Add upload directory to whitelist
]
TEXT_EXTENSIONS = {'.py', '.js', '.json', '.txt', '.md', '.log', '.env', '.yaml', '.yml', '.cfg', '.conf', '.sh', '.html', '.css', '.csv'}
PDF_EXTENSIONS = {'.pdf'}
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp'}
DOCX_EXTENSIONS = {'.docx'}
XLSX_EXTENSIONS = {'.xlsx', '.xls'}
PPTX_EXTENSIONS = {'.pptx', '.ppt'}
CONTEXTUAL_FILE_REFERENCES = {"ini", "itu", "tadi", "tersebut", "file ini", "file itu", "dokumen ini", "dokumen itu"}

# ============================================
# FILE LISTING UTILITY (Reality Check)
# ============================================
def list_my_files(directory: str = None) -> str:
    """
    List all files in a directory with physical verification.
    Default: uploaded_files directory.
    This function ALWAYS checks the actual filesystem - no memory/Mem0 reliance.
    """
    target_dir = directory if directory else UPLOAD_DIR
    
    # Reality check: verify path exists physically
    if not os.path.exists(target_dir):
        return f"Master, folder {target_dir} does not exist physically on disk."
    
    if not os.path.isdir(target_dir):
        return f"Master, {target_dir} is not a directory."
    
    try:
        files = []
        for root, dirs, filenames in os.walk(target_dir):
            for filename in filenames:
                if not filename.startswith('.'):  # Skip hidden files
                    filepath = os.path.join(root, filename)
                    try:
                        stat = os.stat(filepath)
                        size_mb = stat.st_size / (1024 * 1024)
                        mtime = datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M')
                        ext = os.path.splitext(filename)[1].lower()
                        files.append({
                            "name": filename,
                            "path": filepath,
                            "size": f"{size_mb:.2f} MB",
                            "modified": mtime,
                            "type": ext
                        })
                    except Exception:
                        pass
        
        if not files:
            return f"Master, folder {target_dir} is empty. No files found."
        
        # Format output
        result = [f"📁 Files in {target_dir} ({len(files)} files):", ""]
        for f in files:
            icon = "📄"
            if f['type'] in PDF_EXTENSIONS:
                icon = "📕"
            elif f['type'] in IMAGE_EXTENSIONS:
                icon = "🖼️"
            elif f['type'] in {'.py', '.js', '.json'}:
                icon = "💻"
            
            result.append(f"{icon} {f['name']}")
            result.append(f"   Path: {f['path']}")
            result.append(f"   Size: {f['size']} | Modified: {f['modified']}")
            result.append("")
        
        return "\n".join(result)
        
    except PermissionError:
        return f"Master, akses ke folder {target_dir} ditolak. Mohon periksa izin folder."
    except Exception as e:
        return f"Master, error saat membaca folder {target_dir}: {e}"

def list_project_files(directory: str = "/home/kuro/projects/kuro") -> str:
    """
    List all project files in the Kuro project directory.
    Designed for IT Support persona to analyze code structure and detect issues.
    
    Returns:
        Formatted string with file listing organized by type
    """
    target_dir = directory
    
    if not os.path.exists(target_dir):
        return f"Master, project folder {target_dir} does not exist."
    
    if not os.path.isdir(target_dir):
        return f"Master, {target_dir} is not a directory."
    
    try:
        files_by_type = {}
        total_files = 0
        total_size = 0
        
        for root, dirs, filenames in os.walk(target_dir):
            # Skip hidden directories and common non-project folders
            dirs[:] = [d for d in dirs if not d.startswith('.') and d not in ['venv', '__pycache__', 'node_modules', '.git']]
            
            for filename in filenames:
                if filename.startswith('.'):
                    continue
                    
                filepath = os.path.join(root, filename)
                try:
                    stat = os.stat(filepath)
                    size_bytes = stat.st_size
                    total_size += size_bytes
                    total_files += 1
                    
                    ext = os.path.splitext(filename)[1].lower()
                    file_type = ext if ext else 'no-extension'
                    
                    if file_type not in files_by_type:
                        files_by_type[file_type] = []
                    
                    files_by_type[file_type].append({
                        "name": filename,
                        "path": filepath,
                        "size": size_bytes
                    })
                except (PermissionError, OSError):
                    continue
        
        if total_files == 0:
            return f"Master, project folder {target_dir} is empty."
        
        # Format output
        result = [f"📁 Kuro Project Files ({total_files} files, {total_size / (1024*1024):.2f} MB):", ""]
        
        # Sort by file type and show
        for file_type in sorted(files_by_type.keys()):
            files = files_by_type[file_type]
            icon = "📄"
            if file_type == '.py':
                icon = "🐍"
            elif file_type in ['.js', '.ts']:
                icon = "⚡"
            elif file_type in ['.html', '.htm']:
                icon = "🌐"
            elif file_type == '.css':
                icon = "🎨"
            elif file_type == '.json':
                icon = "📋"
            elif file_type == '.md':
                icon = "📝"
            elif file_type == '.log':
                icon = "📜"
            elif file_type in ['.db', '.sqlite']:
                icon = "🗄️"
            
            result.append(f"{icon} {file_type or 'no-ext'} ({len(files)} files)")
            for f in files[:10]:  # Limit per type to avoid overwhelming output
                size_str = f"{f['size'] / 1024:.1f}KB" if f['size'] < 1024*1024 else f"{f['size'] / (1024*1024):.2f}MB"
                result.append(f"   - {f['name']} ({size_str})")
            if len(files) > 10:
                result.append(f"   ... and {len(files) - 10} more")
            result.append("")
        
        return "\n".join(result)
        
    except PermissionError:
        return f"Master, akses ke project folder ditolak. Mohon periksa izin folder."
    except Exception as e:
        return f"Master, error saat membaca project folder: {e}"

# --- Proxmox API Helper ---
def _get_proxmox_headers():
    """Returns Authorization headers for Proxmox API requests."""
    return {"Authorization": f"PVEAPIToken={settings.PVE_TOKEN_ID}={settings.PVE_TOKEN_SECRET}"}


# ============================================
# PDF ENGINE (pdfplumber - Robust with Table Support)
# ============================================

# PDF Processing Timeout Configuration (AFC - Automatic Function Calling)
PDF_PROCESSING_TIMEOUT_SECONDS = 120  # 2 minutes timeout for large PDFs
PDF_CHUNK_PROCESSING_TIMEOUT = 30  # 30 seconds per chunk

def read_pdf_content(
    file_path: str,
    max_pages: int = 20,
    max_chars: int = 15000,
    progress_callback=None  # NEW: Optional callback for SSE "thinking" signals
) -> Dict:
    """
    Robust PDF text extractor using pdfplumber.
    Preserves table structure in Markdown format when possible.
    
    Args:
        file_path: Path to the PDF file
        max_pages: Maximum pages to extract (RAM protection)
        max_chars: Maximum characters to return
        progress_callback: Optional callable(current_page, total_pages) for SSE signals
    
    Returns:
        Dict with 'content', 'page_count', 'tables_found', 'error' keys
    """
    result = {
        "path": file_path,
        "content": "",
        "page_count": 0,
        "tables_found": 0,
        "format": "pdf",
        "error": None
    }
    
    # Reality check: verify file exists physically
    if not os.path.exists(file_path):
        result["error"] = f"File not found: {file_path}"
        return result
    
    try:
        import pdfplumber
    except ImportError:
        result["error"] = "pdfplumber not installed. Run: pip install pdfplumber"
        return result
    
    try:
        with pdfplumber.open(file_path) as pdf:
            result["page_count"] = len(pdf.pages)
            text_parts = []
            total_tables = 0
            total_pages_to_process = min(len(pdf.pages), max_pages)
            
            for i, page in enumerate(pdf.pages[:max_pages]):
                # Send progress signal for SSE "Kuro is thinking..." updates
                if progress_callback:
                    try:
                        progress_callback(current_page=i + 1, total_pages=total_pages_to_process)
                    except Exception as cb_err:
                        logger.debug(f"[PDF_PROGRESS] Callback error (non-critical): {cb_err}")
                
                page_text = f"\n--- Page {i+1} ---\n"
                
                # Extract tables first and convert to Markdown
                tables = page.extract_tables()
                if tables:
                    total_tables += len(tables)
                    for t_idx, table in enumerate(tables):
                        page_text += f"\n**Table {t_idx+1}:**\n"
                        if table and len(table) > 0:
                            # Markdown table format
                            header = table[0]
                            page_text += "| " + " | ".join(str(h or "") for h in header) + " |\n"
                            page_text += "| " + " | ".join("---" for _ in header) + " |\n"
                            for row in table[1:]:
                                page_text += "| " + " | ".join(str(c or "") for c in row) + " |\n"
                        page_text += "\n"
                
                # Extract regular text
                regular_text = page.extract_text()
                if regular_text:
                    page_text += regular_text
                
                text_parts.append(page_text)
            
            result["tables_found"] = total_tables
            result["content"] = "\n".join(text_parts)[:max_chars]
            
            if not result["content"].strip():
                result["error"] = "No text content could be extracted from PDF"
        
        return result
        
    except PermissionError:
        result["error"] = f"Master, akses ke file {file_path} ditolak."
        return result
    except Exception as e:
        result["error"] = f"PDF read error: {e}"
        logger.exception(f"Error reading PDF {file_path}: {e}")
        return result


# ============================================
# UNIVERSAL FILE PARSER ENGINE
# ============================================
def universal_read(file_path: str, max_chars: int = 5000) -> Dict:
    """
    Universal file reader supporting multiple formats.
    
    Args:
        file_path: Path to the file
        max_chars: Maximum characters to return (RAM protection)
    
    Returns:
        Dict with 'content', 'format', 'size', 'error' keys
    """
    result = {"path": file_path, "content": "", "format": "unknown", "size": 0, "error": None}
    
    try:
        if not os.path.exists(file_path):
            result["error"] = f"File not found: {file_path}"
            return result
        
        file_size = os.path.getsize(file_path)
        result["size"] = file_size
        
        # RAM protection: skip large files
        if file_size > MAX_FILE_SIZE_MB * 1024 * 1024:
            result["error"] = f"File too large ({file_size / (1024*1024):.1f}MB). Max: {MAX_FILE_SIZE_MB}MB"
            return result
        
        ext = os.path.splitext(file_path)[1].lower()
        
        # Text/Code files
        if ext in TEXT_EXTENSIONS or ext == '':
            result["format"] = "text"
            try:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    result["content"] = f.read(max_chars)
            except PermissionError:
                result["error"] = f"Master, akses ke file {file_path} ditolak. Mohon periksa izin file."
                logger.warning(result["error"])
        
        # PDF files - use pdfplumber engine
        elif ext in PDF_EXTENSIONS:
            result["format"] = "pdf"
            pdf_result = read_pdf_content(file_path, max_pages=10, max_chars=max_chars)
            if pdf_result.get("error"):
                result["error"] = pdf_result["error"]
            else:
                result["content"] = pdf_result["content"]
                result["tables_found"] = pdf_result.get("tables_found", 0)
                result["page_count"] = pdf_result.get("page_count", 0)
        
        # Image files (return path for Vision API processing)
        elif ext in IMAGE_EXTENSIONS:
            result["format"] = "image"
            result["content"] = f"[Image file: {file_path} - Use Vision API to analyze]"
        
        # Gzipped log files
        elif ext == '.gz':
            result["format"] = "gzip"
            try:
                with gzip.open(file_path, 'rt', encoding='utf-8', errors='replace') as f:
                    result["content"] = f.read(max_chars)
            except PermissionError:
                result["error"] = f"Master, akses ke file {file_path} ditolak."
            except Exception as e:
                result["error"] = f"Gzip read error: {e}"
        
        # Fallback: try as text
        else:
            result["format"] = "binary"
            result["content"] = f"[Binary file: {ext} - {file_size} bytes]"
        
        return result
    
    except PermissionError:
        result["error"] = f"Master, akses ke file {file_path} ditolak. Mohon periksa izin file."
        logger.warning(result["error"])
        return result
    except Exception as e:
        result["error"] = f"Read error: {e}"
        logger.exception(f"Error reading {file_path}: {e}")
        return result


def _run_instruction_on_text(
    instruction: str,
    extracted_text: str,
    source_label: str,
    max_chars: int = 40000,
) -> Dict:
    """Process extracted text with Gemini based on user instruction."""
    from kuro_backend.core import client
    from google.genai import types
    from kuro_backend.config import settings

    trimmed_text = (extracted_text or "")[:max_chars]
    if not trimmed_text.strip():
        return {"success": False, "error": "File tidak mengandung teks yang bisa diproses."}

    prompt = f"""{instruction}

Berikut adalah konten dari {source_label}:
{trimmed_text}

Berikan jawaban terstruktur dan jelas sesuai instruksi."""

    try:
        response = client.models.generate_content(
            model=settings.MODEL_NAME,
            contents=prompt,
            config=types.GenerateContentConfig(temperature=0.2),
        )
        if hasattr(response, "prompt_feedback") and response.prompt_feedback:
            block_reason = getattr(response.prompt_feedback, "block_reason", "UNKNOWN")
            return {"success": False, "error": f"Content blocked by safety filter: {block_reason}"}
        return {"success": True, "output": response.text if response.text else trimmed_text}
    except Exception as exc:
        logger.error(f"Failed to process extracted text with instruction: {exc}")
        return {"success": False, "error": f"Gagal memproses instruksi: {exc}"}


def _extract_image_text_with_vision(image_path: str, instruction: str = "") -> Dict:
    """Extract text/insights from image using Gemini vision."""
    from kuro_backend.core import client
    from google.genai import types
    from kuro_backend.config import settings

    if not os.path.exists(image_path):
        return {"success": False, "error": f"File not found: {image_path}"}

    ext = os.path.splitext(image_path)[1].lower()
    mime_map = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
        ".bmp": "image/bmp",
    }
    mime_type = mime_map.get(ext, "image/jpeg")
    effective_instruction = instruction or "Ekstrak semua teks yang terlihat dari gambar ini secara akurat."

    try:
        with open(image_path, "rb") as f:
            image_data = base64.b64encode(f.read()).decode("utf-8")

        response = client.models.generate_content(
            model=settings.MODEL_NAME,
            contents=[
                types.Part(text=effective_instruction),
                types.Part(
                    inline_data=types.Blob(
                        mime_type=mime_type,
                        data=image_data,
                    )
                ),
            ],
            config=types.GenerateContentConfig(temperature=0.0),
        )
        if hasattr(response, "prompt_feedback") and response.prompt_feedback:
            block_reason = getattr(response.prompt_feedback, "block_reason", "UNKNOWN")
            return {"success": False, "error": f"OCR blocked by safety filter: {block_reason}"}
        return {
            "success": True,
            "content": response.text or "",
            "format": "image_ocr",
            "ocr_engine": "gemini_vision",
        }
    except Exception as exc:
        logger.error(f"Image OCR failed for {image_path}: {exc}")
        return {"success": False, "error": f"Gagal OCR gambar: {exc}"}


def _resolve_smart_read_path(file_ref: str) -> Dict:
    """Resolve explicit or contextual file references."""
    from kuro_backend import memory_manager

    requested = (file_ref or "").strip()
    normalized = requested.lower()
    if not requested or normalized in CONTEXTUAL_FILE_REFERENCES:
        last_accessed = memory_manager.get_runtime_context_value("last_accessed_file", "")
        if not last_accessed:
            return {
                "ok": False,
                "error": "Context unresolved: sebutkan nama/path file dulu agar saya bisa melanjutkan.",
                "resolved_by": "context_missing",
            }
        return {"ok": True, "path": last_accessed, "resolved_by": "context"}

    explicit_path = os.path.abspath(requested)
    if os.path.exists(explicit_path):
        return {"ok": True, "path": explicit_path, "resolved_by": "explicit"}

    upload_candidate = os.path.join(UPLOAD_DIR, requested)
    if os.path.exists(upload_candidate):
        return {"ok": True, "path": upload_candidate, "resolved_by": "search"}

    for root, _, files in os.walk(UPLOAD_DIR):
        for filename in files:
            if requested.lower() in filename.lower():
                return {
                    "ok": True,
                    "path": os.path.join(root, filename),
                    "resolved_by": "search",
                }

    return {
        "ok": False,
        "error": f"File not found: {requested}",
        "resolved_by": "not_found",
    }


def smart_read(file_ref: str = "", instruction: str = "rangkum dokumen ini", max_chars: int = 15000) -> Dict:
    """
    Unified file reading facade with context resolution.
    Supports Office, PDF, image OCR, and plain text/log/code files.
    """
    from kuro_backend import memory_manager

    resolved = _resolve_smart_read_path(file_ref)
    if not resolved.get("ok"):
        logger.warning(
            "[SMART_READ] unresolved: resolved_by=%s file_ref=%s",
            resolved.get("resolved_by"),
            file_ref,
        )
        return {"success": False, "error": resolved.get("error"), "resolved_by": resolved.get("resolved_by")}

    file_path = resolved["path"]
    ext = os.path.splitext(file_path)[1].lower()
    resolved_by = resolved.get("resolved_by", "explicit")
    logger.info("[SMART_READ] start: path=%s resolved_by=%s", file_path, resolved_by)

    try:
        if ext in PDF_EXTENSIONS:
            if os.path.abspath(file_path).startswith(os.path.abspath(UPLOAD_DIR)):
                result = summarize_pdf(os.path.basename(file_path), instruction=instruction)
                engine_used = "summarize_pdf"
            else:
                extracted = read_pdf_content(file_path, max_pages=50, max_chars=max_chars)
                if extracted.get("error"):
                    return {"success": False, "error": extracted["error"], "fallback_reason": "pdf_parse_failed"}
                processed = _run_instruction_on_text(instruction, extracted.get("content", ""), f"PDF file {os.path.basename(file_path)}")
                if not processed.get("success"):
                    return processed
                result = {
                    "success": True,
                    "filename": os.path.basename(file_path),
                    "file_type": "PDF",
                    "summary": processed["output"],
                }
                engine_used = "read_pdf_content+llm"
        elif ext in DOCX_EXTENSIONS:
            extracted = read_docx_content(file_path, max_chars=max_chars)
            if extracted.get("error"):
                return {"success": False, "error": extracted["error"], "fallback_reason": "dependency_missing_or_parse_failed"}
            processed = _run_instruction_on_text(instruction, extracted.get("content", ""), f"Word file {os.path.basename(file_path)}")
            if not processed.get("success"):
                return processed
            result = {
                "success": True,
                "filename": os.path.basename(file_path),
                "file_type": "Word (.docx)",
                "summary": processed["output"],
            }
            engine_used = "read_docx_content+llm"
        elif ext in XLSX_EXTENSIONS:
            extracted = read_xlsx_content(file_path, max_chars=max_chars)
            if extracted.get("error"):
                return {"success": False, "error": extracted["error"], "fallback_reason": "dependency_missing_or_parse_failed"}
            processed = _run_instruction_on_text(instruction, extracted.get("content", ""), f"Excel file {os.path.basename(file_path)}")
            if not processed.get("success"):
                return processed
            result = {
                "success": True,
                "filename": os.path.basename(file_path),
                "file_type": "Excel (.xlsx)",
                "summary": processed["output"],
            }
            engine_used = "read_xlsx_content+llm"
        elif ext in PPTX_EXTENSIONS:
            extracted = read_pptx_content(file_path, max_chars=max_chars)
            if extracted.get("error"):
                return {"success": False, "error": extracted["error"], "fallback_reason": "dependency_missing_or_parse_failed"}
            processed = _run_instruction_on_text(instruction, extracted.get("content", ""), f"PowerPoint file {os.path.basename(file_path)}")
            if not processed.get("success"):
                return processed
            result = {
                "success": True,
                "filename": os.path.basename(file_path),
                "file_type": "PowerPoint (.pptx)",
                "summary": processed["output"],
            }
            engine_used = "read_pptx_content+llm"
        elif ext in IMAGE_EXTENSIONS:
            ocr_result = _extract_image_text_with_vision(file_path, instruction=instruction)
            if not ocr_result.get("success"):
                return ocr_result
            result = {
                "success": True,
                "filename": os.path.basename(file_path),
                "file_type": "Image",
                "summary": ocr_result.get("content", ""),
                "ocr_engine": ocr_result.get("ocr_engine"),
            }
            engine_used = "vision_ocr"
        else:
            raw_result = universal_read(file_path, max_chars=max_chars)
            if raw_result.get("error"):
                return {"success": False, "error": raw_result["error"], "fallback_reason": "universal_read_failed"}
            result = {
                "success": True,
                "filename": os.path.basename(file_path),
                "file_type": raw_result.get("format", "text"),
                "summary": raw_result.get("content", ""),
            }
            engine_used = "universal_read"

        if result.get("success"):
            memory_manager.set_runtime_context_value("last_accessed_file", file_path)
        result["resolved_by"] = resolved_by
        result["engine_used"] = engine_used
        result["source_path"] = file_path
        logger.info(
            "[SMART_READ] done: resolved_by=%s file_type=%s engine_used=%s",
            resolved_by,
            ext or "no_ext",
            engine_used,
        )
        return result
    except Exception as exc:
        logger.exception("[SMART_READ] unexpected error: %s", exc)
        return {
            "success": False,
            "error": f"Smart read failed: {exc}",
            "resolved_by": resolved_by,
            "fallback_reason": "unexpected_exception",
        }


def parse_log_content(content: str) -> List[Dict]:
    """Parse log content using regex for standard log formats."""
    entries = []
    # Common log format: timestamp level message
    log_pattern = re.compile(
        r'(\d{4}-\d{2}-\d{2}[T ]\d{2}:\d{2}:\d{2})[.\d]*\s+(\w+)\s+(.*)',
        re.MULTILINE
    )
    for match in log_pattern.finditer(content):
        entries.append({
            "timestamp": match.group(1),
            "level": match.group(2),
            "message": match.group(3)
        })
    return entries


# ============================================
# RECURSIVE SYSTEM CRAWLER
# ============================================
def index_system_path(path: str, max_files: int = 100) -> Dict:
    """
    Recursively crawl a directory and index files.
    
    Args:
        path: Directory to crawl
        max_files: Maximum files to index (CPU/RAM protection)
    
    Returns:
        Dict with file tree and metadata
    """
    result = {
        "path": path,
        "files": [],
        "total_size": 0,
        "file_count": 0,
        "errors": [],
        "status": "success"
    }
    
    if not os.path.isdir(path):
        result["status"] = "error"
        result["errors"].append(f"Path is not a directory: {path}")
        return result
    
    try:
        for root, dirs, files in os.walk(path):
            if result["file_count"] >= max_files:
                break
            
            # Skip hidden directories and common non-essential paths
            dirs[:] = [d for d in dirs if not d.startswith('.') and d not in {'node_modules', '__pycache__', '.git', 'venv'}]
            
            for filename in files:
                if result["file_count"] >= max_files:
                    break
                
                filepath = os.path.join(root, filename)
                
                try:
                    stat = os.stat(filepath)
                    file_size = stat.st_size
                    
                    # Skip large files
                    if file_size > MAX_FILE_SIZE_MB * 1024 * 1024:
                        continue
                    
                    ext = os.path.splitext(filename)[1].lower()
                    
                    # Read snippet for text files
                    snippet = ""
                    if ext in TEXT_EXTENSIONS and file_size < 1024 * 1024:  # < 1MB
                        try:
                            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                                snippet = f.read(500)
                        except (PermissionError, UnicodeDecodeError):
                            pass
                    
                    result["files"].append({
                        "path": filepath,
                        "name": filename,
                        "size": file_size,
                        "mtime": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                        "extension": ext,
                        "snippet": snippet[:200] if snippet else ""
                    })
                    result["total_size"] += file_size
                    result["file_count"] += 1
                    
                except PermissionError:
                    result["errors"].append(f"Permission denied: {filepath}")
                except Exception as e:
                    result["errors"].append(f"Error indexing {filepath}: {e}")
    
    except PermissionError:
        result["status"] = "error"
        result["errors"].append(f"Master, akses ke direktori {path} ditolak.")
    
    return result


# ============================================
# OS SELF-ANALYSIS LOGIC
# ============================================
def analyze_system_health() -> str:
    """
    Analyze system health by reading /var/log files.
    Reports: failed logins, errors, disk warnings, etc.
    """
    report = ["=== Kuro System Health Analysis ===\n"]
    
    # 1. Auth log analysis (failed logins)
    auth_log = "/var/log/auth.log"
    failed_logins = 0
    sudo_events = 0
    
    try:
        if os.path.exists(auth_log):
            with open(auth_log, 'r', errors='replace') as f:
                content = f.read()
                failed_logins = len(re.findall(r'Failed password', content))
                sudo_events = len(re.findall(r'sudo:', content))
                
            report.append(f"📋 Authentication Log ({auth_log}):")
            report.append(f"   - Failed login attempts: {failed_logins}")
            report.append(f"   - Sudo events: {sudo_events}")
            
            if failed_logins > 3:
                report.append(f"   ⚠️ WARNING: {failed_logins} failed login attempts detected!")
    except PermissionError:
        report.append(f"   ⚠️ Cannot read {auth_log} - permission denied")
    
    # 2. Syslog analysis
    syslog_path = "/var/log/syslog"
    error_count = 0
    
    try:
        if os.path.exists(syslog_path):
            with open(syslog_path, 'r', errors='replace') as f:
                content = f.read()
                error_count = len(re.findall(r'\berror\b|\bERROR\b|\bError\b', content))
            
            report.append(f"\n📋 System Log ({syslog_path}):")
            report.append(f"   - Error entries: {error_count}")
    except PermissionError:
        report.append(f"   ⚠️ Cannot read {syslog_path} - permission denied")
    
    # 3. Disk usage
    report.append(f"\n💾 Disk Usage:")
    for partition in psutil.disk_partitions(all=False):
        try:
            usage = psutil.disk_usage(partition.mountpoint)
            report.append(f"   - {partition.mountpoint}: {usage.percent}% used ({usage.used / (1024**3):.1f}GB / {usage.total / (1024**3):.1f}GB)")
            if usage.percent > 90:
                report.append(f"     ⚠️ CRITICAL: Disk usage above 90%!")
        except PermissionError:
            pass
    
    # 4. Memory status
    ram = psutil.virtual_memory()
    report.append(f"\n🧠 Memory Status:")
    report.append(f"   - Usage: {ram.percent}%")
    report.append(f"   - Available: {ram.available / (1024**3):.2f}GB / {ram.total / (1024**3):.2f}GB")
    
    # 5. Process count
    process_count = len(psutil.pids())
    report.append(f"\n⚙️ Processes: {process_count} running")
    
    # 6. Uptime
    try:
        uptime = subprocess.check_output(['uptime', '-p'], timeout=5).decode('utf-8').strip()
        report.append(f"⏱️ Uptime: {uptime}")
    except Exception:
        pass
    
    return "\n".join(report)


# ============================================
# EXISTING FUNCTIONS (Preserved)
# ============================================
def get_system_status():
    """Mendapatkan status real-time CPU, RAM, Disk, dan Uptime."""
    try:
        cpu_usage = psutil.cpu_percent(interval=0)
        ram = psutil.virtual_memory()
        ram_usage = f"{ram.used / (1024**3):.2f}GB / {ram.total / (1024**3):.2f}GB ({ram.percent}%)"
        disk = psutil.disk_usage('/home/kuro/')
        disk_usage = f"{disk.used / (1024**3):.2f}GB / {disk.total / (1024**3):.2f}GB ({disk.percent}%)"
        uptime = subprocess.check_output(['uptime', '-p'], timeout=5).decode('utf-8').strip()

        return f"Kuro System Health Report:\n- CPU Usage: {cpu_usage}%\n- RAM Usage: {ram_usage}\n- Disk Space (/home/kuro/): {disk_usage}\n- System Uptime: {uptime}"
    except subprocess.TimeoutExpired:
        return "Kuro System Health Report:\n- Uptime: Unable to retrieve (timeout)"
    except Exception as e:
        logger.exception(f"Error in get_system_status: {e}")
        return f"Error retrieving system status: {e}"


def check_proxmox_infrastructure():
    """Mengambil status real-time VM dan Container dari server Proxmox Master."""
    host = settings.PVE_HOST
    headers = _get_proxmox_headers()

    try:
        url = f"https://{host}:8006/api2/json/cluster/resources"
        resp = requests.get(url, headers=headers, verify=False, timeout=10)

        if resp.status_code != 200:
            return f"Error: Proxmox API returned status {resp.status_code}. Response: {resp.text}"

        data = resp.json().get('data', [])
        report = ["Proxmox Infrastructure Audit:"]
        warning_issued = False

        nodes = {}
        for item in data:
            if item['type'] in ['qemu', 'lxc']:
                node = item.get('node', 'unknown')
                if node not in nodes:
                    nodes[node] = []
                nodes[node].append(item)

        for node_name, items in nodes.items():
            report.append(f"\nNode: {node_name}")
            for item in items:
                item_type = item['type'].upper()
                vmid = item.get('vmid', 'unknown')
                name = item.get('name', 'unnamed')
                status = item.get('status', 'unknown')
                report.append(f"  - {item_type} {vmid} ({name}): {status}")
                if status != 'running':
                    report.append(f"    -> WARNING: {item_type} is not running!")
                    warning_issued = True

        if not warning_issued:
            report.append("\nAll systems are running as expected.")

        return "\n".join(report)

    except requests.exceptions.Timeout:
        return f"Error: Connection to Proxmox ({host}:8006) timed out."
    except requests.exceptions.ConnectionError as e:
        return f"Error: Cannot connect to Proxmox ({host}:8006). Detail: {e}"
    except Exception as e:
        logger.exception(f"Error connecting to Proxmox for audit: {e}")
        return f"Error connecting to Proxmox for audit: {e}"


def process_video(video_path: str):
    """Processes a video file using moviepy with RAM safeguards."""
    MAX_VIDEO_SIZE_MB = 100

    try:
        if not os.path.exists(video_path):
            return f"Error: Video file not found: {video_path}"

        file_size_mb = os.path.getsize(video_path) / (1024 * 1024)
        if file_size_mb > MAX_VIDEO_SIZE_MB:
            return f"Error: Video file too large ({file_size_mb:.1f}MB). Maximum: {MAX_VIDEO_SIZE_MB}MB."

        import moviepy as mp
        clip = None
        try:
            clip = mp.VideoFileClip(video_path)
            return f"Video duration: {clip.duration:.2f} seconds."
        finally:
            if clip is not None:
                clip.close()
                del clip

    except ImportError:
        return "Error: moviepy not installed."
    except Exception as e:
        logger.exception(f"Error processing video: {e}")
        return f"Error processing video: {e}"


# ============================================
# REMINDER & SCHEDULING TOOLS
# ============================================
def parse_datetime(text: str) -> Optional[datetime]:
    """
    Parse natural language datetime text into a datetime object.
    Supports: "jam 2 siang", "besok jam 10", "nanti malam", "10 menit lagi", etc.
    Uses settings.TIMEZONE for timezone awareness.
    """
    import pytz
    from kuro_backend.config import settings
    
    tz = settings.tz
    now = datetime.now(tz)
    text_lower = text.lower().strip()
    
    # Relative time patterns
    relative_patterns = {
        "10 menit lagi": lambda: now + timedelta(minutes=10),
        "15 menit lagi": lambda: now + timedelta(minutes=15),
        "30 menit lagi": lambda: now + timedelta(minutes=30),
        "1 jam lagi": lambda: now + timedelta(hours=1),
        "2 jam lagi": lambda: now + timedelta(hours=2),
        "besok": lambda: now + timedelta(days=1),
        "lusa": lambda: now + timedelta(days=2),
        "nanti pagi": lambda: now.replace(hour=8, minute=0, second=0) if now.hour < 8 else (now + timedelta(days=1)).replace(hour=8, minute=0, second=0),
        "nanti siang": lambda: now.replace(hour=12, minute=0, second=0) if now.hour < 12 else (now + timedelta(days=1)).replace(hour=12, minute=0, second=0),
        "nanti sore": lambda: now.replace(hour=16, minute=0, second=0) if now.hour < 16 else (now + timedelta(days=1)).replace(hour=16, minute=0, second=0),
        "nanti malam": lambda: now.replace(hour=20, minute=0, second=0) if now.hour < 20 else (now + timedelta(days=1)).replace(hour=20, minute=0, second=0),
        "hari ini": lambda: now,
    }
    
    for pattern, func in relative_patterns.items():
        if pattern in text_lower:
            return func()
    
    # Time-only patterns: "jam 2 siang", "jam 14:00", "pukul 10"
    time_patterns = [
        r'jam\s+(\d{1,2})[:.](\d{2})',  # jam 14.00
        r'jam\s+(\d{1,2})\s*(pagi|siang|sore|malam)?',  # jam 2 siang
        r'pukul\s+(\d{1,2})[:.](\d{2})',  # pukul 14.00
        r'pukul\s+(\d{1,2})\s*(pagi|siang|sore|malam)?',  # pukul 2 siang
    ]
    
    for pattern in time_patterns:
        match = re.search(pattern, text_lower)
        if match:
            hour = int(match.group(1))
            minute = int(match.group(2)) if match.lastindex >= 2 and match.group(2) else 0
            period = match.group(2) if match.lastindex >= 2 and match.group(2) in ['pagi', 'siang', 'sore', 'malam'] else None
            
            # Adjust hour based on period
            if period:
                if period == 'pagi' and hour < 12:
                    pass  # Already correct
                elif period == 'siang' and hour < 12:
                    hour += 12 if hour != 12 else 0
                elif period == 'sore' and hour < 12:
                    hour += 12
                elif period == 'malam' and hour < 12:
                    hour += 12
            
            target = now.replace(hour=hour, minute=minute, second=0, microsecond=0)
            if target < now:
                target += timedelta(days=1)
            return target
    
    # Full datetime: "2026-04-06 14:00"
    datetime_patterns = [
        r'(\d{4}-\d{2}-\d{2})\s+(\d{1,2})[:.](\d{2})',
        r'(\d{2}/\d{2}/\d{4})\s+(\d{1,2})[:.](\d{2})',
    ]
    
    for pattern in datetime_patterns:
        match = re.search(pattern, text_lower)
        if match:
            try:
                date_str = match.group(1)
                hour = int(match.group(2))
                minute = int(match.group(3))
                
                # Try different date formats
                for fmt in ['%Y-%m-%d', '%d/%m/%Y']:
                    try:
                        dt = datetime.strptime(date_str, fmt)
                        target = dt.replace(hour=hour, minute=minute, tzinfo=tz)
                        return target
                    except ValueError:
                        continue
            except Exception:
                pass
    
    return None


def lookup_chroma_context(query: str) -> str:
    """
    Look up context from Mem0 long-term memory for a reminder query.
    Returns relevant context to enrich the reminder description.
    """
    try:
        from kuro_backend.memory_manager import search_long_term
        results = search_long_term(query, top_k=3)
        if results:
            return "\n".join(results)[:1000]  # Limit context length
    except Exception as e:
        logger.warning(f"Mem0 context lookup failed: {e}")
    return ""


def add_reminder_tool(event_name: str, datetime_text: str, description: str = "", source: str = "web") -> Dict:
    return {"success": False, "error": "Reminders moved to OpenClaw Skills in KURO V7.0"}


def get_reminders_tool() -> Dict:
    return {"success": False, "error": "Reminders moved to OpenClaw Skills in KURO V7.0"}


# ============================================
# FINANCES SSoT (The Chancellor)
# ============================================


def set_monthly_budget_tool(
    month: str,
    amount_usd: float,
    notes: str = "",
) -> Dict[str, Any]:
    """Set or replace the Master's monthly budget (USD) for YYYY-MM in the ledger."""
    from kuro_backend import finance_db

    m = (month or "").strip()
    if len(m) != 7 or m[4] != "-":
        return {"success": False, "error": "month must be YYYY-MM"}
    rid = finance_db.add_budget(m, float(amount_usd), notes or "")
    return {"success": True, "budget_id": rid, "month": m, "amount_usd": float(amount_usd)}


def get_budget_tool(month: str = "") -> Dict[str, Any]:
    """Read the monthly budget for YYYY-MM (defaults to current month)."""
    from datetime import date

    from kuro_backend import finance_db

    m = (month or "").strip()
    if not m:
        m = date.today().strftime("%Y-%m")
    row = finance_db.get_budget(m)
    if not row:
        return {"success": True, "month": m, "budget": None}
    return {"success": True, "month": m, "budget": dict(row)}


def add_recurring_expense_tool(
    label: str,
    amount_usd: float,
    cadence: str = "monthly",
    next_due: str = "",
    category: str = "",
) -> Dict[str, Any]:
    """Record a recurring obligation (subscription, bill) in the finances ledger."""
    from kuro_backend import finance_db

    if not (label or "").strip():
        return {"success": False, "error": "label is required"}
    rid = finance_db.upsert_recurring_expense(
        label.strip(),
        float(amount_usd),
        cadence=cadence or "monthly",
        next_due=next_due or "",
        category=category or "",
        active=True,
    )
    return {"success": True, "expense_id": rid, "label": label.strip()}


def list_recurring_expenses_tool() -> Dict[str, Any]:
    """List active recurring expenses from the finances ledger."""
    from kuro_backend import finance_db

    rows = finance_db.list_recurring_expenses(active_only=True)
    return {"success": True, "expenses": rows}


def get_daily_api_cost_tool(date: str = "") -> Dict[str, Any]:
    """Return estimated API spend (USD) for a calendar day (YYYY-MM-DD); default today."""
    from datetime import date

    from kuro_backend import finance_db

    d = (date or "").strip()
    if not d:
        d = date.today().isoformat()
    cost = finance_db.get_daily_api_cost_usd(d)
    return {"success": True, "date": d, "cost_usd": cost}


# ============================================
# MARKET SENTINEL (OpenClaw readonly)
# ============================================


def _openclaw_skill_body(skill_name: str, payload: Dict[str, Any]) -> Dict[str, Any]:
    from kuro_backend.execution.service import execute_openclaw_skill_sync

    merged = {"execution_mode": "readonly", **(payload or {})}
    return execute_openclaw_skill_sync(skill_name=skill_name, payload=merged)


def get_ticker_price_tool(symbol: str) -> Dict[str, Any]:
    """Fetch last price for a US equity symbol via OpenClaw ``market_analysis`` (readonly)."""
    sym = (symbol or "").strip().upper()
    if not sym:
        return {"success": False, "error": "symbol is required"}
    exec_result = _openclaw_skill_body(
        "market_analysis",
        {"op": "get_ticker_price", "symbol": sym},
    )
    raw = exec_result.get("result") or exec_result.get("raw_response") or {}
    if not isinstance(raw, dict):
        raw = {}
    if not exec_result.get("success"):
        return {
            "success": False,
            "error": exec_result.get("error") or "OpenClaw bridge failure",
            "openclaw": True,
        }
    if raw.get("ok") is False:
        return {
            "success": False,
            "error": raw.get("user_message") or raw.get("error_code") or "market_analysis failed",
            "openclaw": True,
        }
    try:
        from kuro_backend import finance_db

        px = float(raw.get("price"))
        finance_db.apply_watched_price(sym, px)
    except Exception:
        pass
    return {"success": True, "openclaw": True, **raw}


def get_market_news_tool(symbol: str) -> Dict[str, Any]:
    """Headlines for a symbol via OpenClaw ``market_analysis`` get_market_news (may be empty)."""
    sym = (symbol or "").strip().upper()
    if not sym:
        return {"success": False, "error": "symbol is required"}
    exec_result = _openclaw_skill_body(
        "market_analysis",
        {"op": "get_market_news", "symbol": sym},
    )
    raw = exec_result.get("result") or exec_result.get("raw_response") or {}
    if not isinstance(raw, dict):
        raw = {}
    if not exec_result.get("success"):
        return {
            "success": False,
            "error": exec_result.get("error") or "OpenClaw bridge failure",
            "openclaw": True,
        }
    if raw.get("ok") is False:
        return {
            "success": False,
            "error": raw.get("user_message") or "market_analysis failed",
            "openclaw": True,
        }
    return {"success": True, "openclaw": True, **raw}


def prediction_market_scan_tool(topics: str = "") -> Dict[str, Any]:
    """Prediction-style probability rows from OpenClaw ``prediction_market_scan`` (readonly)."""
    topic_list: List[str] = []
    raw_t = (topics or "").strip()
    if raw_t:
        topic_list = [t.strip() for t in raw_t.split(",") if t.strip()]
    exec_result = _openclaw_skill_body(
        "prediction_market_scan",
        {"topics": topic_list} if topic_list else {},
    )
    raw = exec_result.get("result") or exec_result.get("raw_response") or {}
    if not isinstance(raw, dict):
        raw = {}
    if not exec_result.get("success"):
        return {
            "success": False,
            "error": exec_result.get("error") or "OpenClaw bridge failure",
            "openclaw": True,
        }
    if raw.get("ok") is False:
        return {
            "success": False,
            "error": raw.get("user_message") or "prediction_market_scan failed",
            "openclaw": True,
        }
    markets = raw.get("markets") or []
    if isinstance(markets, list):
        try:
            from kuro_backend import finance_db

            existing = {
                str(r["slug"]): float(r.get("last_probability") or 0.0)
                for r in finance_db.list_prediction_watch()
            }
            for m in markets:
                if not isinstance(m, dict):
                    continue
                slug = str(m.get("topic_id") or "").strip()
                if not slug:
                    continue
                prob = float(m.get("probability") or 0.0)
                prev = existing.get(slug)
                trend = "flat"
                if prev is not None:
                    if prob > prev + 0.005:
                        trend = "up"
                    elif prob < prev - 0.005:
                        trend = "down"
                finance_db.upsert_prediction_watch(
                    slug,
                    str(m.get("title") or slug)[:240],
                    prob,
                    trend=trend,
                )
                existing[slug] = prob
        except Exception as exc:
            logger.debug("[MARKET] prediction cache upsert skipped: %s", exc)
    return {"success": True, "openclaw": True, **raw}


# ============================================
# DAILY HABIT TOOLS
# ============================================
def mark_habit_done_tool(habit_title: str) -> Dict:
    """
    Mark a daily habit as done via natural language.
    E.g., "Aku udah gym ya hari ini" -> marks "Gym" as done.
    """
    return {"success": False, "error": "Habits moved to OpenClaw Skills in KURO V7.0"}


def get_habits_status_tool() -> Dict:
    return {"success": False, "error": "Habits moved to OpenClaw Skills in KURO V7.0"}


EMPTY_HABIT_FACTUAL_MESSAGE = "Saya tidak menemukan catatan data faktual."


def get_habit_history_tool(habit_title: str, days: int = 30) -> Dict:
    """
    Factual habit log rows from SQLite for one habit. If has_factual_data is false, you MUST
    reply to the user with exactly: 'Saya tidak menemukan catatan data faktual.' — no ISO clauses,
    no IP addresses, no invented activities.
    """
    return {
        "success": False,
        "has_factual_data": False,
        "ai_required_reply": "Habits moved to OpenClaw Skills in KURO V7.0",
        "habit_logs": [],
        "completion_dates": [],
    }


# ============================================
# PDF SUMMARIZATION TOOL
# ============================================
def summarize_pdf(pdf_filename: str, instruction: str = "rangkum dokumen ini") -> Dict:
    """
    Full PDF workflow: Find file -> Read PDF -> Chunk if needed -> Send to Gemini -> Return summary.
    
    Args:
        pdf_filename: Name of the PDF file (e.g., "VCT26 Official Competition Ruleset.pdf")
        instruction: What to do with the PDF content (e.g., "rangkum", "jelaskan poin penting")
    
    Returns:
        Dict with summary or error message
    """
    from kuro_backend.core import client
    from google.genai import types
    from kuro_backend.config import settings
    
    # Step 1: Find the file in uploaded_files
    pdf_path = os.path.join(UPLOAD_DIR, pdf_filename)
    
    # Also try searching in uploaded_files recursively
    if not os.path.exists(pdf_path):
        for root, dirs, files in os.walk(UPLOAD_DIR):
            for f in files:
                if pdf_filename.lower() in f.lower():
                    pdf_path = os.path.join(root, f)
                    break
            if os.path.exists(pdf_path):
                break
    
    if not os.path.exists(pdf_path):
        return {
            "success": False,
            "error": f"Master, file '{pdf_filename}' tidak ditemukan di folder uploaded_files."
        }
    
    # Step 2: Read and extract text from PDF
    pdf_result = read_pdf_content(pdf_path, max_pages=50, max_chars=50000)
    
    if pdf_result.get("error"):
        return {
            "success": False,
            "error": f"Gagal membaca PDF: {pdf_result['error']}"
        }
    
    extracted_text = pdf_result.get("content", "")
    
    if not extracted_text.strip():
        return {
            "success": False,
            "error": "PDF tidak mengandung teks yang bisa diekstrak. Mungkin PDF berbasis gambar (scanned)."
        }
    
    # Step 3: Handle large text with chunking if needed
    MAX_TOKENS_APPROX = 30000  # Approximate character limit for Gemini
    text_length = len(extracted_text)
    
    if text_length > MAX_TOKENS_APPROX:
        # Chunk the text
        chunks = []
        chunk_size = MAX_TOKENS_APPROX // 2
        for i in range(0, text_length, chunk_size):
            chunks.append(extracted_text[i:i + chunk_size])
        
        logger.info(f"PDF text chunked into {len(chunks)} parts for processing.")
        
        # Process each chunk and combine summaries
        chunk_summaries = []
        for idx, chunk in enumerate(chunks):
            chunk_prompt = f"""{instruction}

Bagian {idx+1}/{len(chunks)} dari dokumen:
{chunk}

Berikan rangkuman untuk bagian ini saja."""
            
            try:
                response = client.models.generate_content(
                    model=settings.MODEL_NAME,
                    contents=chunk_prompt,
                    config=types.GenerateContentConfig(temperature=0.2)
                )
                # SAFETY CHECK
                if hasattr(response, 'prompt_feedback') and response.prompt_feedback:
                    logger.warning(f"[SUMMARIZE_PDF] Chunk blocked: {getattr(response.prompt_feedback, 'block_reason', 'UNKNOWN')}")
                    continue
                if response.text:
                    chunk_summaries.append(response.text)
            except Exception as e:
                logger.error(f"Error processing chunk {idx+1}: {e}")
        
        # Combine chunk summaries into final summary
        if chunk_summaries:
            combined_text = "\n\n".join(chunk_summaries)
            final_prompt = f"""{instruction}

Berikut adalah rangkuman dari setiap bagian dokumen:
{combined_text}

Berikan rangkuman final yang menyeluruh dan terstruktur."""
            
            try:
                final_response = client.models.generate_content(
                    model=settings.MODEL_NAME,
                    contents=final_prompt,
                    config=types.GenerateContentConfig(temperature=0.2)
                )
                # SAFETY CHECK
                if hasattr(final_response, 'prompt_feedback') and final_response.prompt_feedback:
                    logger.warning(f"[SUMMARIZE_PDF] Final response blocked: {getattr(final_response.prompt_feedback, 'block_reason', 'UNKNOWN')}")
                    return {"success": False, "error": "Content blocked by safety filter"}
                extracted_text = final_response.text if final_response.text else combined_text
            except Exception as e:
                logger.error(f"Error creating final summary: {e}")
                extracted_text = combined_text
        else:
            return {
                "success": False,
                "error": "Gagal memproses semua chunk PDF."
            }
    else:
        # Step 4: Send to Gemini for processing (single chunk)
        full_prompt = f"""{instruction}

Berikut adalah isi dokumen PDF "{pdf_filename}":
{extracted_text}

Berikan hasil yang diminta dengan format yang rapi dan terstruktur."""
        
        try:
            response = client.models.generate_content(
                model=settings.MODEL_NAME,
                contents=full_prompt,
                config=types.GenerateContentConfig(temperature=0.2)
            )
            # SAFETY CHECK
            if hasattr(response, 'prompt_feedback') and response.prompt_feedback:
                block_reason = getattr(response.prompt_feedback, 'block_reason', 'UNKNOWN')
                logger.warning(f"[SUMMARIZE_PDF] Content blocked: {block_reason}")
                return {"success": False, "error": f"Content blocked by safety filter: {block_reason}"}
            extracted_text = response.text if response.text else extracted_text
        except Exception as e:
            logger.error(f"Error sending to Gemini: {e}")
            return {
                "success": False,
                "error": f"Gagal memproses dengan Gemini: {e}"
            }
    
    # Step 5: Return the result
    return {
        "success": True,
        "filename": pdf_filename,
        "pages": pdf_result.get("page_count", 0),
        "tables_found": pdf_result.get("tables_found", 0),
        "summary": extracted_text
    }


# ============================================
# WORD DOCUMENT (.docx) EXTRACTION
# ============================================
def read_docx_content(file_path: str, max_chars: int = 15000) -> Dict:
    """
    Extract text from Word document (.docx) including paragraphs and tables.
    
    Args:
        file_path: Path to the .docx file
        max_chars: Maximum characters to extract
    
    Returns:
        Dict with content, page_count, tables_found, or error
    """
    try:
        from docx import Document
        
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        
        doc = Document(file_path)
        text_parts = []
        table_count = 0
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            table_count += 1
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        
        full_text = "\n\n".join(text_parts)
        
        # Truncate if too long
        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[... dokumen dipotong karena terlalu panjang ...]"
        
        return {
            "success": True,
            "content": full_text,
            "char_count": len(full_text),
            "tables_found": table_count,
            "format": "docx"
        }
        
    except ImportError:
        return {"error": "python-docx not installed. Run: pip install python-docx"}
    except Exception as e:
        logger.error(f"Error reading DOCX: {e}")
        return {"error": f"Failed to read Word document: {str(e)}"}


# ============================================
# EXCEL SPREADSHEET (.xlsx) EXTRACTION
# ============================================
def read_xlsx_content(file_path: str, max_chars: int = 15000) -> Dict:
    """
    Extract data from Excel spreadsheet (.xlsx) as Markdown tables.
    
    Args:
        file_path: Path to the .xlsx file
        max_chars: Maximum characters to extract
    
    Returns:
        Dict with content, sheet_count, or error
    """
    try:
        from openpyxl import load_workbook
        
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        
        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheet_count = len(wb.sheetnames)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"\n## Sheet: {sheet_name}\n")
            
            # Convert sheet to Markdown table
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                text_parts.append("(Sheet kosong)")
                continue
            
            # Find header row (first non-empty row)
            header_idx = 0
            for i, row in enumerate(rows):
                if any(cell is not None for cell in row):
                    header_idx = i
                    break
            
            # Get headers
            headers = [str(cell) if cell is not None else "" for cell in rows[header_idx]]
            
            # Build Markdown table
            header_row = "| " + " | ".join(headers) + " |"
            separator = "| " + " | ".join("---" for _ in headers) + " |"
            text_parts.append(header_row)
            text_parts.append(separator)
            
            # Add data rows
            for row in rows[header_idx + 1:]:
                cells = [str(cell) if cell is not None else "" for cell in row]
                # Skip empty rows
                if not any(cells):
                    continue
                data_row = "| " + " | ".join(cells) + " |"
                text_parts.append(data_row)
        
        wb.close()
        
        full_text = "\n".join(text_parts)
        
        # Truncate if too long
        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[... data dipotong karena terlalu panjang ...]"
        
        return {
            "success": True,
            "content": full_text,
            "char_count": len(full_text),
            "sheet_count": sheet_count,
            "format": "xlsx"
        }
        
    except ImportError:
        return {"error": "openpyxl not installed. Run: pip install openpyxl"}
    except Exception as e:
        logger.error(f"Error reading XLSX: {e}")
        return {"error": f"Failed to read Excel file: {str(e)}"}


# ============================================
# POWERPOINT (.pptx) EXTRACTION
# ============================================
def read_pptx_content(file_path: str, max_chars: int = 15000) -> Dict:
    """
    Extract text from PowerPoint presentation (.pptx) including slides and notes.
    
    Args:
        file_path: Path to the .pptx file
        max_chars: Maximum characters to extract
    
    Returns:
        Dict with content, slide_count, or error
    """
    try:
        from pptx import Presentation
        
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        text_parts = []
        
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = f"\n### Slide {idx}\n"
            
            # Extract title
            if slide.shapes.title:
                slide_text += f"**{slide.shapes.title.text}**\n\n"
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Skip title (already added)
                    if shape == slide.shapes.title:
                        continue
                    slide_text += shape.text + "\n"
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide and notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_text += f"\n*Catatan: {notes_text}*\n"
            
            text_parts.append(slide_text)
        
        full_text = "\n".join(text_parts)
        
        # Truncate if too long
        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[... presentasi dipotong karena terlalu panjang ...]"
        
        return {
            "success": True,
            "content": full_text,
            "char_count": len(full_text),
            "slide_count": slide_count,
            "format": "pptx"
        }
        
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}
    except Exception as e:
        logger.error(f"Error reading PPTX: {e}")
        return {"error": f"Failed to read PowerPoint file: {str(e)}"}


# ============================================
# UNIVERSAL DOCUMENT SUMMARIZATION TOOL
# ============================================
def summarize_document(filename: str, instruction: str = "rangkum dokumen ini") -> Dict:
    """
    Universal document summarization: Auto-detect file type and extract content.
    Supports: PDF, DOCX, XLSX, PPTX
    
    Args:
        filename: Name of the file (e.g., "Laporan.docx", "Data.xlsx")
        instruction: What to do with the content (e.g., "rangkum", "jelaskan poin penting")
    
    Returns:
        Dict with summary or error message
    """
    from kuro_backend.core import client
    from google.genai import types
    from kuro_backend.config import settings
    
    # Step 1: Find the file in uploaded_files
    file_path = os.path.join(UPLOAD_DIR, filename)
    
    # Also try searching recursively
    if not os.path.exists(file_path):
        for root, dirs, files in os.walk(UPLOAD_DIR):
            for f in files:
                if filename.lower() in f.lower():
                    file_path = os.path.join(root, f)
                    break
            if os.path.exists(file_path):
                break
    
    if not os.path.exists(file_path):
        return {
            "success": False,
            "error": f"Master, file '{filename}' tidak ditemukan di folder uploaded_files."
        }
    
    # Step 2: Detect file type and extract content
    ext = os.path.splitext(file_path)[1].lower()
    extracted_text = ""
    file_info = {}
    
    if ext in PDF_EXTENSIONS:
        pdf_result = read_pdf_content(file_path, max_pages=50, max_chars=50000)
        if pdf_result.get("error"):
            return {"success": False, "error": f"Gagal membaca PDF: {pdf_result['error']}"}
        extracted_text = pdf_result.get("content", "")
        file_info = {"type": "PDF", "pages": pdf_result.get("page_count", 0), "tables": pdf_result.get("tables_found", 0)}
    
    elif ext in DOCX_EXTENSIONS:
        docx_result = read_docx_content(file_path, max_chars=50000)
        if docx_result.get("error"):
            return {"success": False, "error": f"Gagal membaca Word: {docx_result['error']}"}
        extracted_text = docx_result.get("content", "")
        file_info = {"type": "Word (.docx)", "tables": docx_result.get("tables_found", 0)}
    
    elif ext in XLSX_EXTENSIONS:
        xlsx_result = read_xlsx_content(file_path, max_chars=50000)
        if xlsx_result.get("error"):
            return {"success": False, "error": f"Gagal membaca Excel: {xlsx_result['error']}"}
        extracted_text = xlsx_result.get("content", "")
        file_info = {"type": "Excel (.xlsx)", "sheets": xlsx_result.get("sheet_count", 0)}
    
    elif ext in PPTX_EXTENSIONS:
        pptx_result = read_pptx_content(file_path, max_chars=50000)
        if pptx_result.get("error"):
            return {"success": False, "error": f"Gagal membaca PowerPoint: {pptx_result['error']}"}
        extracted_text = pptx_result.get("content", "")
        file_info = {"type": "PowerPoint (.pptx)", "slides": pptx_result.get("slide_count", 0)}
    
    else:
        return {
            "success": False,
            "error": f"Master, format file '{ext}' belum didukung. Format yang didukung: PDF, DOCX, XLSX, PPTX"
        }
    
    if not extracted_text.strip():
        return {
            "success": False,
            "error": "File tidak mengandung teks yang bisa diekstrak."
        }
    
    # Step 3: Handle large text with chunking if needed
    MAX_TOKENS_APPROX = 30000
    text_length = len(extracted_text)
    
    if text_length > MAX_TOKENS_APPROX:
        chunks = []
        chunk_size = MAX_TOKENS_APPROX // 2
        for i in range(0, text_length, chunk_size):
            chunks.append(extracted_text[i:i + chunk_size])
        
        logger.info(f"Document text chunked into {len(chunks)} parts for processing.")
        
        chunk_summaries = []
        for idx, chunk in enumerate(chunks):
            chunk_prompt = f"""{instruction}

Bagian {idx+1}/{len(chunks)} dari dokumen:
{chunk}

Berikan rangkuman untuk bagian ini saja."""
            
            try:
                response = client.models.generate_content(
                    model=settings.MODEL_NAME,
                    contents=chunk_prompt,
                    config=types.GenerateContentConfig(temperature=0.2)
                )
                # SAFETY CHECK
                if hasattr(response, 'prompt_feedback') and response.prompt_feedback:
                    logger.warning(f"[SUMMARIZE_DOC] Chunk blocked: {getattr(response.prompt_feedback, 'block_reason', 'UNKNOWN')}")
                    continue
                if response.text:
                    chunk_summaries.append(response.text)
            except Exception as e:
                logger.error(f"Error processing chunk {idx+1}: {e}")
        
        if chunk_summaries:
            combined_text = "\n\n".join(chunk_summaries)
            final_prompt = f"""{instruction}

Berikut adalah rangkuman dari setiap bagian dokumen:
{combined_text}

Berikan rangkuman final yang menyeluruh dan terstruktur."""
            
            try:
                final_response = client.models.generate_content(
                    model=settings.MODEL_NAME,
                    contents=final_prompt,
                    config=types.GenerateContentConfig(temperature=0.2)
                )
                # SAFETY CHECK
                if hasattr(final_response, 'prompt_feedback') and final_response.prompt_feedback:
                    logger.warning(f"[SUMMARIZE_DOC] Final response blocked: {getattr(final_response.prompt_feedback, 'block_reason', 'UNKNOWN')}")
                    return {"success": False, "error": "Content blocked by safety filter"}
                extracted_text = final_response.text if final_response.text else combined_text
            except Exception as e:
                logger.error(f"Error creating final summary: {e}")
                extracted_text = combined_text
        else:
            return {"success": False, "error": "Gagal memproses semua chunk dokumen."}
    else:
        # Step 4: Send to Gemini for processing (single chunk)
        full_prompt = f"""{instruction}

Berikut adalah isi dokumen {file_info.get('type', 'unknown')} "{filename}":
{extracted_text}

Berikan hasil yang diminta dengan format yang rapi dan terstruktur."""
        
        try:
            response = client.models.generate_content(
                model=settings.MODEL_NAME,
                contents=full_prompt,
                config=types.GenerateContentConfig(temperature=0.2)
            )
            # SAFETY CHECK
            if hasattr(response, 'prompt_feedback') and response.prompt_feedback:
                block_reason = getattr(response.prompt_feedback, 'block_reason', 'UNKNOWN')
                logger.warning(f"[SUMMARIZE_DOC] Content blocked: {block_reason}")
                return {"success": False, "error": f"Content blocked by safety filter: {block_reason}"}
            extracted_text = response.text if response.text else extracted_text
        except Exception as e:
            logger.error(f"Error sending to Gemini: {e}")
            return {"success": False, "error": f"Gagal memproses dengan Gemini: {e}"}
    
    # Step 5: Return the result
    return {
        "success": True,
        "filename": filename,
        "file_type": file_info.get("type", "unknown"),
        "summary": extracted_text
    }


def _run_async_coro_sync(coro):
    """
    Run async coroutine from sync context.
    If already inside a running loop, execute in a dedicated worker thread.
    """
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        return asyncio.run(coro)

    result_holder = {"value": None, "error": None}

    def _runner():
        try:
            result_holder["value"] = asyncio.run(coro)
        except Exception as exc:
            result_holder["error"] = exc

    thread = threading.Thread(target=_runner, daemon=True)
    thread.start()
    thread.join()

    if result_holder["error"] is not None:
        raise result_holder["error"]
    return result_holder["value"]


GEMINI_SHARE_URL_PATTERN = re.compile(
    r"https://gemini\.google\.com/share/[^\s)]+", re.IGNORECASE
)
GEMINI_HARVEST_SKILL_NAME = "harvest_gemini_share"
GEMINI_SHARE_BLOCKED_USER_MESSAGE = (
    "Master, akses ke link Gemini ini tersumbat (Timeout/Bot Protection). "
    "Mohon cek kembali atau berikan file PDF-nya saja."
)


def extract_gemini_share_url(text: str) -> Optional[str]:
    """Return first Gemini share URL in text, or None."""
    if not text:
        return None
    m = GEMINI_SHARE_URL_PATTERN.search(text)
    if not m:
        return None
    return m.group(0).rstrip(".,;]")


def task_suggests_gemini_harvest(task: str) -> bool:
    """Heuristic: user wants research library ingest from a Gemini share link."""
    t = (task or "").lower()
    keys = (
        "library riset",
        "masukkan",
        "pelajari",
        "gemini share",
        "link gemini",
        "share ini",
        "riset",
    )
    return any(k in t for k in keys)


def resolve_harvest_gemini_routing(
    task: str,
    skill_name: str,
    params: Optional[Dict],
) -> Tuple[str, Dict]:
    """
    If task contains a Gemini share URL and ingest intent, route to harvest_gemini_share.
    Ensures share_url is present in params when using that skill.
    """
    merged = dict(params or {})
    url = merged.get("share_url") or merged.get("url") or extract_gemini_share_url(task)
    skill = (skill_name or "general_execution").strip() or "general_execution"

    if url and isinstance(url, str):
        url = url.strip()
        if skill in ("general_execution",) and task_suggests_gemini_harvest(task):
            skill = GEMINI_HARVEST_SKILL_NAME
        if skill == GEMINI_HARVEST_SKILL_NAME:
            merged.setdefault("share_url", url)

    return skill, merged


def _openclaw_body_implies_blocked_or_timeout(raw: Dict, execution_error: Optional[str]) -> bool:
    if raw.get("error_code") == "blocked_or_timeout":
        return True
    if raw.get("user_message") == GEMINI_SHARE_BLOCKED_USER_MESSAGE:
        return True
    blob = " ".join(
        str(x).lower()
        for x in (raw.get("detail"), raw.get("error"), execution_error)
        if x
    )
    return "timeout" in blob or "blocked_or_timeout" in blob or "bot" in blob


def advanced_execution_tool(
    task_description: str,
    params: Optional[Dict] = None,
    skill_name: str = "general_execution",
    execution_mode: str = "mutating",
) -> Dict:
    """
    Delegate advanced operational/system tasks to OpenClaw daemon.

    Args:
        task_description: Natural language execution brief from model/user.
        params: Optional structured parameters.
        skill_name: OpenClaw skill identifier (default: general_execution).
    """
    task = (task_description or "").strip()
    if not task:
        return {
            "success": False,
            "error": "task_description wajib diisi untuk advanced_execution_tool.",
        }

    skill_name, params = resolve_harvest_gemini_routing(task, skill_name, params)

    normalized_mode = "readonly" if str(execution_mode).strip().lower() == "readonly" else "mutating"
    payload = {
        "task_description": task,
        "execution_mode": normalized_mode,
        **(params or {}),
    }

    try:
        from kuro_backend.execution.service import execute_openclaw_skill_sync

        execution_result = execute_openclaw_skill_sync(skill_name=skill_name, payload=payload)
    except Exception as exc:
        logger.exception("[OPENCLAW_TOOL] Bridge execution failed: %s", exc)
        err_msg = str(exc).lower()
        if skill_name == GEMINI_HARVEST_SKILL_NAME and (
            "timeout" in err_msg or "timed out" in err_msg
        ):
            return {
                "success": False,
                "error": GEMINI_SHARE_BLOCKED_USER_MESSAGE,
                "message": GEMINI_SHARE_BLOCKED_USER_MESSAGE,
                "skill_name": skill_name,
                "task_description": task,
            }
        return {
            "success": False,
            "error": f"OpenClaw bridge error: {exc}",
            "message": "Eksekusi OpenClaw gagal dijalankan.",
        }

    success = bool(execution_result.get("success"))
    raw = execution_result.get("result") or execution_result.get("raw_response") or {}
    if not isinstance(raw, dict):
        raw = {"raw": raw}

    # Daemon may return HTTP 200 with ok=false in JSON body
    if raw.get("ok") is False:
        success = False

    bridge_err = execution_result.get("error")
    if skill_name == GEMINI_HARVEST_SKILL_NAME and not success:
        if _openclaw_body_implies_blocked_or_timeout(raw, bridge_err if isinstance(bridge_err, str) else None):
            execution_result = dict(execution_result)
            execution_result["error"] = GEMINI_SHARE_BLOCKED_USER_MESSAGE

    from kuro_backend import memory_coordinator as _mem_coord

    bump_meta = _mem_coord.apply_openclaw_execution_result(
        success=success, skill_name=skill_name, raw=raw
    )
    should_bump_revision = bool(bump_meta.get("should_bump_revision"))
    revision_bumped = bool(bump_meta.get("revision_bumped"))
    revision_bump_error = bump_meta.get("revision_error")

    message = "Tugas berhasil didelegasikan ke OpenClaw."
    memory_fallback_required = bool(execution_result.get("memory_fallback_required"))
    memory_fallback_saved = False
    memory_fallback_error = None

    if memory_fallback_required:
        try:
            from kuro_backend import memory_manager

            fallback_note = (
                "[OPENCLAW_UNAVAILABLE_FALLBACK]\n"
                f"skill_name={skill_name}\n"
                f"task_description={task}\n"
                f"params={json.dumps(params or {}, ensure_ascii=False)}"
            )
            memory_manager.add_short_term("system", fallback_note)
            memory_fallback_saved = True
        except Exception as exc:
            memory_fallback_error = str(exc)
            logger.exception("[OPENCLAW_TOOL] Failed to persist memory fallback: %s", exc)

    if not success:
        message = "Eksekusi OpenClaw gagal."
        if skill_name == GEMINI_HARVEST_SKILL_NAME and _openclaw_body_implies_blocked_or_timeout(
            raw, bridge_err if isinstance(bridge_err, str) else None
        ):
            message = GEMINI_SHARE_BLOCKED_USER_MESSAGE
        if memory_fallback_required:
            message = execution_result.get("error") or message
            if memory_fallback_saved:
                message += " Instruksi sudah dicatat ke memori sementara."
            elif memory_fallback_error:
                message += " Namun pencatatan memori sementara gagal."
    elif revision_bumped:
        message += " Data revision berhasil diperbarui untuk sinkronisasi dashboard."
    elif should_bump_revision and revision_bump_error:
        message += " Namun bump data revision gagal."

    response = {
        "success": success,
        "message": message,
        "execution_result": execution_result,
        "skill_name": skill_name,
        "task_description": task,
        "ssot_sync": {
            "requested": should_bump_revision,
            "revision_bumped": revision_bumped,
            "error": revision_bump_error,
        },
        "memory_fallback": {
            "required": memory_fallback_required,
            "saved": memory_fallback_saved,
            "error": memory_fallback_error,
        },
    }
    if not success:
        err_out = execution_result.get("error")
        if skill_name == GEMINI_HARVEST_SKILL_NAME and _openclaw_body_implies_blocked_or_timeout(
            raw, err_out if isinstance(err_out, str) else None
        ):
            response["error"] = GEMINI_SHARE_BLOCKED_USER_MESSAGE
        elif err_out:
            response["error"] = err_out
    return response
